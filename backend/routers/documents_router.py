from fastapi import APIRouter, HTTPException, Depends, status
from fastapi.responses import FileResponse, StreamingResponse
from firebase_admin import firestore
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
from datetime import datetime
import io
import google.generativeai as genai
import os
from dotenv import load_dotenv
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup

from .auth_router import verify_token

load_dotenv()

router = APIRouter()


def get_db():
  return firestore.client()


gemini_api_key = os.getenv("GEMINI_API_KEY")
if gemini_api_key:
  genai.configure(api_key=gemini_api_key)


class RefinementRequest(BaseModel):
  section_id: str
  prompt: str


class FeedbackRequest(BaseModel):
  section_id: str
  feedback_type: str
  comment: Optional[str] = None


class CommentRequest(BaseModel):
  section_id: str
  comment: str


class UpdateSectionRequest(BaseModel):
  section_id: str
  content: str


def add_html_to_docx(paragraph, html_content):
  soup = BeautifulSoup(html_content, "html.parser")

  # If no tags, treat as plain text
  if not soup.find():
    paragraph.add_run(html_content)
    return

  for child in soup.descendants:
    if child.name in ["strong", "b"]:
      run = paragraph.add_run(child.get_text())
      run.bold = True
    elif child.name == "li":
      run = paragraph.add_run(f"\n• {child.get_text()}")
    elif child.name is None:
      if child.parent.name not in ["strong", "b", "li"]:
        paragraph.add_run(child)


def generate_content_with_gemini(
  topic: str,
  section_title: str,
  document_type: str,
  context: Optional[str] = None,
) -> str:
  if not gemini_api_key:
    return f"""
Placeholder content for {section_title}.
"""

  models_to_try = ["gemini-2.5-flash", "gemini-1.5-flash", "gemini-pro"]
  for model_name in models_to_try:
    try:
      model = genai.GenerativeModel(model_name)

      if document_type == "docx":
        prompt = f"""
Act as a professional technical writer. Write a comprehensive section titled "{section_title}" for a document about "{topic}".

Context: {context if context else "Start of document"}

STRICT OUTPUT RULES:
1. Use HTML tags for formatting (<p>, <ul>, <li>).
2. Do NOT use Markdown.
3. Length: 300 words.
4. Do NOT include the title.
"""
      else:
        prompt = f"""
Act as a presentation expert. Write content for a slide titled "{section_title}" for a deck about "{topic}".

Context: {context if context else "Start"}

STRICT OUTPUT RULES:
1. Use HTML tags (<p>, <ul>, <li>).
2. Use bullet points.
3. Max 100 words.
4. Do NOT include the title.
"""

      response = model.generate_content(prompt)
      return (
        response.text.replace("```html", "").replace("```", "")
      )
    except Exception:
      continue

  return """
Error generating content.
"""


def generate_template_with_gemini(
  topic: str, document_type: str
) -> List[Dict[str, Any]]:
  if not gemini_api_key:
    return []

  models_to_try = ["gemini-2.5-flash", "gemini-1.5-flash", "gemini-pro"]

  for model_name in models_to_try:
    try:
      model = genai.GenerativeModel(model_name)
      prompt = (
        f"Generate a {'5-8 section' if document_type == 'docx' else '8-10 slide'} "
        f"outline for a {'document' if document_type == 'docx' else 'presentation'} "
        f"about '{topic}'.\nReturn ONLY titles, one per line."
      )
      response = model.generate_content(prompt)
      titles = [
        line.strip().replace("*", "").replace("-", "").strip()
        for line in response.text.split("\n")
        if line.strip()
      ]
      return [
        {"title": t, "order": i} for i, t in enumerate(titles)
      ]
    except Exception:
      continue

  return []


@router.post("/{project_id}/generate-template")
async def generate_template(
  project_id: str, decoded_token: dict = Depends(verify_token)
):
  user_id = decoded_token["uid"]
  doc_ref = get_db().collection("projects").document(project_id)
  doc = doc_ref.get()

  if not doc.exists or doc.to_dict()["user_id"] != user_id:
    raise HTTPException(status_code=404)

  data = doc.to_dict()
  template = generate_template_with_gemini(
    data["topic"], data["document_type"]
  )
  return {"template": template}


@router.post("/{project_id}/generate-content")
async def generate_content(
  project_id: str, decoded_token: dict = Depends(verify_token)
):
  user_id = decoded_token["uid"]
  db = get_db()
  doc_ref = db.collection("projects").document(project_id)
  doc = doc_ref.get()

  if not doc.exists or doc.to_dict()["user_id"] != user_id:
    raise HTTPException(status_code=404)

  project_data = doc.to_dict()
  items = (
    project_data.get("outline", [])
    if project_data["document_type"] == "docx"
    else project_data.get("slides", [])
  )
  content_map = project_data.get("content", {})
  context = ""
  updated_items: List[Dict[str, Any]] = []
  has_changes = False

  for item in sorted(items, key=lambda x: x.get("order", 0)):
    if not item.get("id") or str(item.get("id")).startswith("temp"):
      item["id"] = (
        f"item_{int(datetime.now().timestamp() * 1000)}_"
        f"{item.get('order')}"
      )
      has_changes = True

    if (
      item["id"] not in content_map
      or not content_map[item["id"]].get("content")
    ):
      gen_text = generate_content_with_gemini(
        project_data["topic"],
        item["title"],
        project_data["document_type"],
        context,
      )
      content_map[item["id"]] = {
        "title": item["title"],
        "content": gen_text,
        "order": item["order"],
      }

    raw_text = BeautifulSoup(
      content_map[item["id"]]["content"], "html.parser"
    ).get_text()
    context += f"\n{item['title']}: {raw_text[:150]}..."
    updated_items.append(item)

  update_payload: Dict[str, Any] = {
    "content": content_map,
    "updated_at": datetime.utcnow().isoformat(),
  }

  if has_changes:
    key = (
      "outline"
      if project_data["document_type"] == "docx"
      else "slides"
    )
    update_payload[key] = updated_items

  doc_ref.update(update_payload)
  return {"content": content_map}


@router.post("/{project_id}/refine")
async def refine_content(
  project_id: str,
  refinement: RefinementRequest,
  decoded_token: dict = Depends(verify_token),
):
  user_id = decoded_token["uid"]
  db = get_db()
  doc_ref = db.collection("projects").document(project_id)
  doc = doc_ref.get()

  if not doc.exists or doc.to_dict()["user_id"] != user_id:
    raise HTTPException(status_code=404)

  project_data = doc.to_dict()
  content = project_data.get("content", {})

  if refinement.section_id not in content:
    raise HTTPException(status_code=404)

  current_content = content[refinement.section_id]["content"]
  models_to_try = ["gemini-2.5-flash", "gemini-1.5-flash", "gemini-pro"]

  for model_name in models_to_try:
    try:
      model = genai.GenerativeModel(model_name)
      prompt = f"""
Refine this text (HTML format): {current_content}
Instruction: {refinement.prompt}

IMPORTANT: Return valid HTML (<p>, <ul>, <li>).
"""
      response = model.generate_content(prompt)
      refined_text = (
        response.text.replace("```html", "").replace("```", "")
      )
      content[refinement.section_id]["content"] = refined_text
      break
    except Exception:
      continue

  doc_ref.update(
    {
      "content": content,
      "refinement_history": firestore.ArrayUnion(
        [
          {
            "type": "ai_refine",
            "section_id": refinement.section_id,
            "value": refinement.prompt,
            "timestamp": datetime.utcnow().isoformat(),
          }
        ]
      ),
    }
  )

  return {
    "refined_content": content[refinement.section_id]["content"]
  }


@router.get("/{project_id}/export")
async def export_document(
  project_id: str, decoded_token: dict = Depends(verify_token)
):
  user_id = decoded_token["uid"]
  db = get_db()
  doc_ref = db.collection("projects").document(project_id)
  doc = doc_ref.get()

  if not doc.exists or doc.to_dict()["user_id"] != user_id:
    raise HTTPException(status_code=404)

  project_data = doc.to_dict()
  content_map = project_data.get("content", {})
  items = (
    project_data.get("outline", [])
    if project_data["document_type"] == "docx"
    else project_data.get("slides", [])
  )
  sorted_items = sorted(items, key=lambda x: x.get("order", 0))

  file_stream = io.BytesIO()

  def get_html(item):
    if item.get("id") in content_map:
      return content_map[item["id"]].get("content", "")
    for _, v in content_map.items():
      if v.get("order") == item.get("order"):
        return v.get("content", "")
    return ""

  # DOCX
  if project_data["document_type"] == "docx":
    docx_doc = Document()
    docx_doc.add_heading(project_data.get("topic", "Document"), 0)

    for item in sorted_items:
      docx_doc.add_heading(item.get("title", "Untitled"), level=1)
      html_content = get_html(item)
      soup = BeautifulSoup(html_content, "html.parser")

      paragraphs = soup.find_all("p")
      if not paragraphs:
        p = docx_doc.add_paragraph()
        add_html_to_docx(p, html_content)
      else:
        for para_tag in paragraphs:
          p = docx_doc.add_paragraph()
          add_html_to_docx(p, str(para_tag))

      ul_tags = soup.find_all("ul")
      for ul in ul_tags:
        for li in ul.find_all("li"):
          p = docx_doc.add_paragraph(style="List Bullet")
          add_html_to_docx(p, str(li.decode_contents()))

    docx_doc.save(file_stream)
    media_type = (
      "application/vnd.openxmlformats-officedocument."
      "wordprocessingml.document"
    )
    filename = f"{project_data.get('name')}.docx"

  # PPTX
  else:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = project_data.get("name")
    slide.placeholders[1].text = project_data.get("topic")

    for item in sorted_items:
      slide = prs.slides.add_slide(prs.slide_layouts[1])
      slide.shapes.title.text = item.get("title")

      html_content = get_html(item)
      soup = BeautifulSoup(html_content, "html.parser")
      text_frame = slide.placeholders[1].text_frame
      text_frame.clear()
      p = text_frame.paragraphs[0]
      p.text = soup.get_text("\n")

    prs.save(file_stream)
    media_type = (
      "application/vnd.openxmlformats-officedocument."
      "presentationml.presentation"
    )
    filename = f"{project_data.get('name')}.pptx"

  file_stream.seek(0)
  return StreamingResponse(
    file_stream,
    media_type=media_type,
    headers={"Content-Disposition": f'attachment; filename="{filename}"'},
  )


# --- Feedback / Comments / Manual edits history ---


@router.post("/{project_id}/feedback")
async def add_feedback(
  project_id: str,
  feedback: FeedbackRequest,
  decoded_token: dict = Depends(verify_token),
):
  user_id = decoded_token["uid"]
  db = get_db()
  doc_ref = db.collection("projects").document(project_id)

  history_entry = {
    "type": "feedback",
    "section_id": feedback.section_id,
    "value": feedback.feedback_type,
    "timestamp": datetime.utcnow().isoformat(),
  }

  doc_ref.update(
    {
      "refinement_history": firestore.ArrayUnion([history_entry])
    }
  )
  return {"message": "ok"}


@router.post("/{project_id}/comment")
async def add_comment(
  project_id: str,
  comment: CommentRequest,
  decoded_token: dict = Depends(verify_token),
):
  user_id = decoded_token["uid"]
  db = get_db()
  doc_ref = db.collection("projects").document(project_id)

  history_entry = {
    "type": "comment",
    "section_id": comment.section_id,
    "value": comment.comment,
    "timestamp": datetime.utcnow().isoformat(),
  }

  doc_ref.update(
    {
      "refinement_history": firestore.ArrayUnion([history_entry])
    }
  )
  return {"message": "ok"}


@router.post("/{project_id}/update-section")
async def update_section(
  project_id: str,
  update: UpdateSectionRequest,
  decoded_token: dict = Depends(verify_token),
):
  """Persist manual edits from the rich-text editor for a single section."""
  user_id = decoded_token["uid"]
  db = get_db()
  doc_ref = db.collection("projects").document(project_id)
  doc = doc_ref.get()

  if not doc.exists or doc.to_dict()["user_id"] != user_id:
    raise HTTPException(status_code=404)

  project_data = doc.to_dict()
  content = project_data.get("content", {})

  if update.section_id not in content:
    raise HTTPException(
      status_code=404, detail="Section not found"
    )

  content[update.section_id]["content"] = update.content

  history_entry = {
    "type": "manual_edit",
    "section_id": update.section_id,
    "value": "manual edit",
    "timestamp": datetime.utcnow().isoformat(),
  }

  doc_ref.update(
    {
      "content": content,
      "updated_at": datetime.utcnow().isoformat(),
      "refinement_history": firestore.ArrayUnion([history_entry]),
    }
  )

  return {"message": "ok"}
